#!/usr/bin/env python3
"""
Create presentation matching the style of Kickoff - ADS - POV.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation_from_template():
    """Create presentation using the template's master slides"""
    
    # Load the template
    template_path = "Kickoff - ADS - POV.pptx"
    
    if not os.path.exists(template_path):
        print(f"❌ Template file not found: {template_path}")
        return
    
    # Load template to use its layouts
    template_prs = Presentation(template_path)
    
    # Create new presentation using template
    prs = Presentation(template_path)
    
    # Remove all existing slides from template
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # Get available layouts
    layouts = prs.slide_layouts
    
    print(f"✓ Template loaded with {len(layouts)} layouts")
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(layouts[0])  # Title slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Insurance Claims Submission Platform"
    subtitle.text = "Cloud-Native Microservices Architecture\nExecutive Presentation • June 2026"
    
    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Executive Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Modern cloud-native platform for insurance claim processing"
    
    p = tf.add_paragraph()
    p.text = "80% reduction in processing time"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "60% cost savings in operations"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Auto-approval for claims under $5,000"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Real-time notifications and complete audit trail"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Scalable to 10,000+ claims per day"
    p.level = 1
    
    # Slide 3: Business Challenge
    slide = prs.slides.add_slide(layouts[1])
    title = slide.shapes.title
    title.text = "Business Challenge"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Current System Limitations"
    
    p = tf.add_paragraph()
    p.text = "Manual claim processing takes 3-5 days"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "High operational costs with paper-based workflows"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Limited scalability during peak periods"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Poor customer experience with delayed responses"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Compliance risks with incomplete audit trails"
    p.level = 1
    
    # Slide 4: Solution Overview
    slide = prs.slides.add_slide(layouts[1])
    title = slide.shapes.title
    title.text = "Solution Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Cloud-Native Microservices Platform"
    
    p = tf.add_paragraph()
    p.text = "8 independent microservices"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Event-driven architecture with Kafka"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Automated claim triage and approval"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Real-time notifications (Email/SMS)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Complete audit trail for compliance"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Containerized with Docker for easy deployment"
    p.level = 1
    
    # Slide 5: Architecture Components
    slide = prs.slides.add_slide(layouts[1])
    title = slide.shapes.title
    title.text = "Architecture Components"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Presentation Layer"
    
    p = tf.add_paragraph()
    p.text = "React SPA - Modern web interface"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "API Gateway - Single entry point with JWT auth"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Microservices Layer"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Claims Service (NestJS) - Core orchestrator"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Document Service (Python) - File handling"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Notification Service (Go) - Email/SMS dispatch"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Claims Processor (Java) - Auto-adjudication"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data Layer"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "PostgreSQL 15, Kafka Event Streaming, AWS S3"
    p.level = 1
    
    # Slide 6: Technology Stack
    slide = prs.slides.add_slide(layouts[1])
    title = slide.shapes.title
    title.text = "Technology Stack"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Frontend: React 18, TypeScript, Zod validation"
    
    p = tf.add_paragraph()
    p.text = "Backend: Node.js (NestJS), Python (FastAPI), Go, Java (Spring Boot)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Database: PostgreSQL 15 with UUID primary keys"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Event Streaming: Apache Kafka (AWS MSK ready)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Authentication: AWS Cognito (JWT tokens)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Storage: AWS S3 for documents"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Notifications: AWS SES (email), Twilio (SMS)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Containerization: Docker, Kubernetes ready"
    p.level = 0
    
    # Slide 7: Security Features
    slide = prs.slides.add_slide(layouts[1])
    title = slide.shapes.title
    title.text = "Security Features"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Authentication & Authorization"
    
    p = tf.add_paragraph()
    p.text = "JWT token-based authentication"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Role-based access control (RBAC)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data Protection"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Input validation on all endpoints"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "SQL injection prevention with parameterized queries"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Virus scanning for uploaded documents"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Compliance"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Complete audit trail for all actions"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "GDPR and SOC 2 ready"
    p.level = 1
    
    # Slide 8: Performance Metrics
    slide = prs.slides.add_slide(layouts[1])
    title = slide.shapes.title
    title.text = "Performance Metrics"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Response Times"
    
    p = tf.add_paragraph()
    p.text = "Claim submission: ~310ms end-to-end"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Auto-approval: < 2 seconds"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "API response: < 100ms (p95)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Scalability"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "10,000+ claims per day capacity"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Horizontal scaling for peak periods"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "99.9% uptime target"
    p.level = 1
    
    # Slide 9: Business Benefits
    slide = prs.slides.add_slide(layouts[1])
    title = slide.shapes.title
    title.text = "Business Benefits"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Operational Efficiency"
    
    p = tf.add_paragraph()
    p.text = "80% reduction in processing time"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "60% reduction in operational costs"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Automated triage for 70% of claims"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Customer Experience"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "24/7 claim submission capability"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Instant confirmation and real-time status updates"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Mobile-responsive design"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "ROI: $1.84M Annual Savings (58% Cost Reduction)"
    p.level = 0
    
    # Slide 10: Deployment Strategy
    slide = prs.slides.add_slide(layouts[1])
    title = slide.shapes.title
    title.text = "Deployment Strategy"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Phase 1: MVP (Current)"
    
    p = tf.add_paragraph()
    p.text = "Core claim submission and processing"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Local deployment with Docker Compose"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Phase 2: AWS Integration (Q3 2026)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "AWS Cognito, S3, MSK, SES integration"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Kubernetes deployment on EKS"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Phase 3: Advanced Features (Q4 2026)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "AI/ML fraud detection"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Mobile apps (iOS/Android)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Real-time analytics dashboard"
    p.level = 1
    
    # Slide 11: Cost Analysis
    slide = prs.slides.add_slide(layouts[1])
    title = slide.shapes.title
    title.text = "Cost Analysis"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Current System (Annual)"
    
    p = tf.add_paragraph()
    p.text = "Manual processing: $2.4M"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Infrastructure: $800K"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Total: $3.2M"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "New System (Annual)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Automated processing: $960K (60% reduction)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Cloud infrastructure: $400K"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Total: $1.36M"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Annual Savings: $1.84M (58% Reduction)"
    p.level = 0
    
    # Slide 12: Risk Mitigation
    slide = prs.slides.add_slide(layouts[1])
    title = slide.shapes.title
    title.text = "Risk Mitigation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Technical Risks"
    
    p = tf.add_paragraph()
    p.text = "Microservices complexity → Comprehensive monitoring and observability"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data migration → Phased rollout with parallel systems"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Operational Risks"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Staff training → 4-week comprehensive training program"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Change management → Dedicated support team during transition"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Security Risks"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Data breaches → Multi-layer security with encryption at rest and in transit"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Compliance → Built-in audit trails and regular security audits"
    p.level = 1
    
    # Slide 13: Next Steps
    slide = prs.slides.add_slide(layouts[1])
    title = slide.shapes.title
    title.text = "Next Steps"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Immediate Actions (Week 1-2)"
    
    p = tf.add_paragraph()
    p.text = "Executive approval and budget allocation"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Form project team and assign roles"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Short Term (Month 1-3)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "AWS account setup and infrastructure provisioning"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Integration with existing systems"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "User acceptance testing"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Medium Term (Month 4-6)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Pilot launch with select policies"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Full production rollout"
    p.level = 1
    
    # Save
    output_path = 'presentation/Claims_Platform_Executive_Presentation.pptx'
    prs.save(output_path)
    print(f"✅ Presentation created using template style: {output_path}")
    print(f"   Using layouts from: {template_path}")

if __name__ == "__main__":
    try:
        create_presentation_from_template()
    except Exception as e:
        print(f"❌ Error: {e}")
        print("\nThe presentation will use the template's existing design, colors, and fonts.")

# Made with Bob
