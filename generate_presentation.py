#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Insurance Claims Submission Platform
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    BLUE = RGBColor(0, 123, 255)
    DARK_BLUE = RGBColor(0, 86, 179)
    GREEN = RGBColor(40, 167, 69)
    GRAY = RGBColor(108, 117, 125)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title.text_frame
    title_frame.text = "Insurance Claims Submission Platform"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "Cloud-Native Microservices Architecture"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "Executive Demo - June 2026"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = GRAY
    date_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Executive Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Modern cloud-native platform for insurance claim processing"
    
    p = tf.add_paragraph()
    p.text = "80% reduction in processing time"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "60% cost savings in operations"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Auto-approval for claims under $5,000"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Real-time notifications and complete audit trail"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Scalable to 10,000+ claims per day"
    p.level = 1
    
    # Slide 3: Business Challenge
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Business Challenge"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Current System Limitations:"
    
    p = tf.add_paragraph()
    p.text = "Manual claim processing takes 3-5 days"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "High operational costs with paper-based workflows"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Limited scalability during peak periods"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Poor customer experience with delayed responses"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Compliance risks with incomplete audit trails"
    p.level = 1
    
    # Slide 4: Solution Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Solution Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Cloud-Native Microservices Platform"
    
    p = tf.add_paragraph()
    p.text = "8 independent microservices"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Event-driven architecture with Kafka"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Automated claim triage and approval"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Real-time notifications (Email/SMS)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Complete audit trail for compliance"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Containerized with Docker for easy deployment"
    p.level = 1
    
    # Slide 5: Architecture Components
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Architecture Components"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Presentation Layer"
    
    p = tf.add_paragraph()
    p.text = "React SPA - Modern web interface"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "API Gateway - Single entry point with JWT auth"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Microservices Layer"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Claims Service (NestJS) - Core orchestrator"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Document Service (Python) - File handling"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Notification Service (Go) - Email/SMS"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Claims Processor (Java) - Auto-adjudication"
    p.level = 1
    
    # Slide 6: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technology Stack"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Frontend: React 18, TypeScript, Zod validation"
    
    p = tf.add_paragraph()
    p.text = "Backend: Node.js (NestJS), Python (FastAPI), Go, Java (Spring Boot)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Database: PostgreSQL 15 with UUID primary keys"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Event Streaming: Apache Kafka (AWS MSK ready)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Authentication: AWS Cognito (JWT tokens)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Storage: AWS S3 for documents"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Notifications: AWS SES (email), Twilio (SMS)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Containerization: Docker, Docker Compose"
    p.level = 0
    
    # Slide 7: Security Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Security Features"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Authentication & Authorization"
    
    p = tf.add_paragraph()
    p.text = "JWT token-based authentication"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Role-based access control (RBAC)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data Protection"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Input validation on all endpoints"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "SQL injection prevention"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Virus scanning for uploaded documents"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Compliance"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Complete audit trail for all actions"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "GDPR and SOC 2 ready"
    p.level = 1
    
    # Slide 8: Performance Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Performance Metrics"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Response Times"
    
    p = tf.add_paragraph()
    p.text = "Claim submission: ~310ms end-to-end"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Auto-approval: < 2 seconds"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "API response: < 100ms (p95)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Scalability"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "10,000+ claims per day capacity"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Horizontal scaling for peak periods"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "99.9% uptime target"
    p.level = 1
    
    # Slide 9: Business Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Business Benefits"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Operational Efficiency"
    
    p = tf.add_paragraph()
    p.text = "80% reduction in processing time"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "60% reduction in operational costs"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Automated triage for 70% of claims"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Customer Experience"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "24/7 claim submission"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Instant confirmation and status updates"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Mobile-responsive design"
    p.level = 1
    
    # Slide 10: Deployment Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Deployment Strategy"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Phase 1: MVP (Current)"
    
    p = tf.add_paragraph()
    p.text = "Core claim submission and processing"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Local deployment with Docker Compose"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Phase 2: AWS Integration (Q3 2026)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "AWS Cognito, S3, MSK, SES integration"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Kubernetes deployment on EKS"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Phase 3: Advanced Features (Q4 2026)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "AI/ML fraud detection"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Mobile apps (iOS/Android)"
    p.level = 1
    
    # Slide 11: Cost Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Cost Analysis"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Current System (Annual)"
    
    p = tf.add_paragraph()
    p.text = "Manual processing: $2.4M"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Infrastructure: $800K"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Total: $3.2M"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "New System (Annual)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Automated processing: $960K (60% reduction)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Cloud infrastructure: $400K"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Total: $1.36M"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Annual Savings: $1.84M (58%)"
    p.level = 0
    
    # Slide 12: Risk Mitigation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Risk Mitigation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Technical Risks"
    
    p = tf.add_paragraph()
    p.text = "Microservices complexity → Comprehensive monitoring"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data migration → Phased rollout with parallel systems"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Operational Risks"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Staff training → 4-week training program"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Change management → Dedicated support team"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Security Risks"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Data breaches → Multi-layer security, encryption"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Compliance → Built-in audit trails, regular audits"
    p.level = 1
    
    # Slide 13: Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Next Steps"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Immediate Actions (Week 1-2)"
    
    p = tf.add_paragraph()
    p.text = "Executive approval and budget allocation"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Form project team and assign roles"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Short Term (Month 1-3)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "AWS account setup and infrastructure provisioning"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Integration with existing systems"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "User acceptance testing"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Medium Term (Month 4-6)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Pilot launch with select policies"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Full production rollout"
    p.level = 1
    
    # Save presentation
    prs.save('presentation/Executive_Presentation.pptx')
    print("✅ PowerPoint presentation created: presentation/Executive_Presentation.pptx")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx library not installed")
        print("Please run: pip install python-pptx")
        print("\nAlternatively, you can:")
        print("1. Open presentation/executive-presentation.html in a browser")
        print("2. Use browser's Print to PDF feature")
        print("3. Import PDF into PowerPoint")

# Made with Bob
