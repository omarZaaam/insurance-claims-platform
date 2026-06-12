#!/usr/bin/env python3
"""
Generate Styled PowerPoint presentation for Canada Life
Insurance Claims Submission Platform
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_canada_life_header(slide, title_text):
    """Add Canada Life branded header to slide"""
    # Canada Life Green
    CL_GREEN = RGBColor(0, 167, 88)
    
    # Add green header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = CL_GREEN
    header.line.fill.background()
    
    # Add title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.LEFT

def add_footer(slide, slide_number):
    """Add Canada Life footer"""
    CL_GREEN = RGBColor(0, 167, 88)
    
    # Footer bar
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7),
        Inches(10), Inches(0.5)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = CL_GREEN
    footer.line.fill.background()
    
    # Canada Life text
    footer_text = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.1),
        Inches(4), Inches(0.3)
    )
    tf = footer_text.text_frame
    tf.text = "Canada Life"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Page number
    page_num = slide.shapes.add_textbox(
        Inches(9), Inches(7.1),
        Inches(0.5), Inches(0.3)
    )
    tf = page_num.text_frame
    tf.text = str(slide_number)
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.RIGHT

def add_icon_bullet(slide, left, top, icon_text, text, level=0):
    """Add a bullet point with icon"""
    CL_GREEN = RGBColor(0, 167, 88)
    DARK_GRAY = RGBColor(51, 51, 51)
    
    # Icon circle
    icon = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + level * 0.3), Inches(top),
        Inches(0.25), Inches(0.25)
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = CL_GREEN
    icon.line.fill.background()
    
    # Icon text
    icon_tf = icon.text_frame
    icon_tf.text = icon_text
    icon_p = icon_tf.paragraphs[0]
    icon_p.font.size = Pt(12)
    icon_p.font.bold = True
    icon_p.font.color.rgb = RGBColor(255, 255, 255)
    icon_p.alignment = PP_ALIGN.CENTER
    icon_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Text
    text_box = slide.shapes.add_textbox(
        Inches(left + 0.35 + level * 0.3), Inches(top - 0.05),
        Inches(8.5 - level * 0.3), Inches(0.35)
    )
    tf = text_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(16 if level == 0 else 14)
    p.font.color.rgb = DARK_GRAY

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Colors
    CL_GREEN = RGBColor(0, 167, 88)
    CL_DARK_GREEN = RGBColor(0, 120, 63)
    WHITE = RGBColor(255, 255, 255)
    DARK_GRAY = RGBColor(51, 51, 51)
    LIGHT_GRAY = RGBColor(242, 242, 242)
    
    # ===== SLIDE 1: Title =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Full green background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = CL_GREEN
    bg.line.fill.background()
    
    # Canada Life logo placeholder
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.6))
    logo_tf = logo_box.text_frame
    logo_tf.text = "CANADA LIFE"
    logo_p = logo_tf.paragraphs[0]
    logo_p.font.size = Pt(36)
    logo_p.font.bold = True
    logo_p.font.color.rgb = WHITE
    
    # Main title
    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    tf = title.text_frame
    tf.text = "Insurance Claims\nSubmission Platform"
    for para in tf.paragraphs:
        para.font.size = Pt(48)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.6))
    tf = subtitle.text_frame
    tf.text = "Cloud-Native Microservices Architecture"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.4))
    tf = date.text_frame
    tf.text = "Executive Presentation • June 2026"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # ===== SLIDE 2: Executive Summary =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Executive Summary")
    add_footer(slide, 2)
    
    y = 1.2
    add_icon_bullet(slide, 0.8, y, "✓", "Modern cloud-native platform for insurance claim processing")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "⚡", "80% reduction in processing time")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "💰", "60% cost savings in operations")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "🎯", "Auto-approval for claims under $5,000")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "📱", "Real-time notifications and complete audit trail")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "📈", "Scalable to 10,000+ claims per day")
    
    # ===== SLIDE 3: Business Challenge =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Business Challenge")
    add_footer(slide, 3)
    
    # Challenge box
    challenge_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.2),
        Inches(8.4), Inches(1.2)
    )
    challenge_box.fill.solid()
    challenge_box.fill.fore_color.rgb = LIGHT_GRAY
    challenge_box.line.color.rgb = CL_GREEN
    challenge_box.line.width = Pt(2)
    
    challenge_text = slide.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(7.6), Inches(0.8))
    tf = challenge_text.text_frame
    tf.text = "Current System Limitations"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CL_DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    y = 2.8
    add_icon_bullet(slide, 0.8, y, "⏱", "Manual claim processing takes 3-5 days")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "📄", "High operational costs with paper-based workflows")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "📊", "Limited scalability during peak periods")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "😞", "Poor customer experience with delayed responses")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "⚠", "Compliance risks with incomplete audit trails")
    
    # ===== SLIDE 4: Solution Overview =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Solution Overview")
    add_footer(slide, 4)
    
    # Solution highlight box
    solution_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.2),
        Inches(8.4), Inches(1)
    )
    solution_box.fill.solid()
    solution_box.fill.fore_color.rgb = CL_GREEN
    solution_box.line.fill.background()
    
    solution_text = slide.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(7.6), Inches(0.6))
    tf = solution_text.text_frame
    tf.text = "Cloud-Native Microservices Platform"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    y = 2.6
    add_icon_bullet(slide, 0.8, y, "🔧", "8 independent microservices")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "⚡", "Event-driven architecture with Kafka")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "🤖", "Automated claim triage and approval")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "📧", "Real-time notifications (Email/SMS)")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "📝", "Complete audit trail for compliance")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "🐳", "Containerized with Docker for easy deployment")
    
    # ===== SLIDE 5: Architecture Components =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Architecture Components")
    add_footer(slide, 5)
    
    # Three-tier architecture boxes
    # Presentation Layer
    layer1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.3),
        Inches(8), Inches(1.2)
    )
    layer1.fill.solid()
    layer1.fill.fore_color.rgb = RGBColor(230, 247, 237)
    layer1.line.color.rgb = CL_GREEN
    layer1.line.width = Pt(2)
    
    layer1_title = slide.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(7.6), Inches(0.3))
    tf = layer1_title.text_frame
    tf.text = "Presentation Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CL_DARK_GREEN
    
    y = 1.8
    add_icon_bullet(slide, 1.2, y, "⚛", "React SPA - Modern web interface", 1)
    y += 0.4
    add_icon_bullet(slide, 1.2, y, "🚪", "API Gateway - Single entry point with JWT auth", 1)
    
    # Microservices Layer
    layer2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2.8),
        Inches(8), Inches(2.2)
    )
    layer2.fill.solid()
    layer2.fill.fore_color.rgb = RGBColor(230, 247, 237)
    layer2.line.color.rgb = CL_GREEN
    layer2.line.width = Pt(2)
    
    layer2_title = slide.shapes.add_textbox(Inches(1.2), Inches(2.9), Inches(7.6), Inches(0.3))
    tf = layer2_title.text_frame
    tf.text = "Microservices Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CL_DARK_GREEN
    
    y = 3.3
    add_icon_bullet(slide, 1.2, y, "📋", "Claims Service (NestJS) - Core orchestrator", 1)
    y += 0.4
    add_icon_bullet(slide, 1.2, y, "📄", "Document Service (Python) - File handling", 1)
    y += 0.4
    add_icon_bullet(slide, 1.2, y, "📧", "Notification Service (Go) - Email/SMS", 1)
    y += 0.4
    add_icon_bullet(slide, 1.2, y, "⚙", "Claims Processor (Java) - Auto-adjudication", 1)
    
    # Data Layer
    layer3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(5.3),
        Inches(8), Inches(1)
    )
    layer3.fill.solid()
    layer3.fill.fore_color.rgb = RGBColor(230, 247, 237)
    layer3.line.color.rgb = CL_GREEN
    layer3.line.width = Pt(2)
    
    layer3_title = slide.shapes.add_textbox(Inches(1.2), Inches(5.4), Inches(7.6), Inches(0.3))
    tf = layer3_title.text_frame
    tf.text = "Data Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CL_DARK_GREEN
    
    y = 5.8
    add_icon_bullet(slide, 1.2, y, "🗄", "PostgreSQL 15 • Kafka Event Streaming • AWS S3", 1)
    
    # ===== SLIDE 6: Technology Stack =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Technology Stack")
    add_footer(slide, 6)
    
    y = 1.2
    add_icon_bullet(slide, 0.8, y, "⚛", "Frontend: React 18, TypeScript, Zod validation")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "🔧", "Backend: Node.js, Python, Go, Java (Spring Boot)")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "🗄", "Database: PostgreSQL 15 with UUID primary keys")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "📨", "Event Streaming: Apache Kafka (AWS MSK ready)")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "🔐", "Authentication: AWS Cognito (JWT tokens)")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "☁", "Storage: AWS S3 for documents")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "📧", "Notifications: AWS SES (email), Twilio (SMS)")
    y += 0.5
    add_icon_bullet(slide, 0.8, y, "🐳", "Containerization: Docker, Kubernetes ready")
    
    # ===== SLIDE 7: Security Features =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Security Features")
    add_footer(slide, 7)
    
    # Security pillars
    pillar_width = 2.5
    pillar_height = 4.5
    pillar_spacing = 0.3
    
    for i, (title, items) in enumerate([
        ("Authentication", ["🔑 JWT tokens", "👤 RBAC", "🔐 MFA ready"]),
        ("Data Protection", ["✓ Input validation", "🛡 SQL injection prevention", "🦠 Virus scanning"]),
        ("Compliance", ["📝 Audit trail", "📊 GDPR ready", "🏛 SOC 2 ready"])
    ]):
        x = 0.8 + i * (pillar_width + pillar_spacing)
        
        pillar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.5),
            Inches(pillar_width), Inches(pillar_height)
        )
        pillar.fill.solid()
        pillar.fill.fore_color.rgb = LIGHT_GRAY
        pillar.line.color.rgb = CL_GREEN
        pillar.line.width = Pt(2)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(1.7),
            Inches(pillar_width - 0.2), Inches(0.4)
        )
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = CL_DARK_GREEN
        p.alignment = PP_ALIGN.CENTER
        
        # Items
        y_item = 2.3
        for item in items:
            item_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y_item),
                Inches(pillar_width - 0.4), Inches(0.4)
            )
            tf = item_box.text_frame
            tf.text = item
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            y_item += 0.6
    
    # ===== SLIDE 8: Performance Metrics =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Performance Metrics")
    add_footer(slide, 8)
    
    # Metric cards
    metrics = [
        ("⚡", "310ms", "Claim Submission"),
        ("🚀", "< 2 sec", "Auto-Approval"),
        ("📊", "10,000+", "Claims/Day"),
        ("✓", "99.9%", "Uptime Target")
    ]
    
    card_width = 2
    card_height = 2
    spacing = 0.3
    start_x = 1
    start_y = 2
    
    for i, (icon, value, label) in enumerate(metrics):
        row = i // 2
        col = i % 2
        x = start_x + col * (card_width + spacing)
        y = start_y + row * (card_height + spacing)
        
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(card_width), Inches(card_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CL_GREEN
        card.line.fill.background()
        
        # Icon
        icon_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.2),
            Inches(card_width - 0.2), Inches(0.4)
        )
        tf = icon_box.text_frame
        tf.text = icon
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER
        
        # Value
        value_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.7),
            Inches(card_width - 0.2), Inches(0.5)
        )
        tf = value_box.text_frame
        tf.text = value
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 1.3),
            Inches(card_width - 0.2), Inches(0.4)
        )
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Additional metrics
    y = 5
    add_icon_bullet(slide, 5.5, y, "📈", "Horizontal scaling for peak periods")
    y += 0.5
    add_icon_bullet(slide, 5.5, y, "⚡", "API response < 100ms (p95)")
    
    # ===== SLIDE 9: Business Benefits =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Business Benefits")
    add_footer(slide, 9)
    
    # Two columns
    # Left column - Operational
    col1_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.3),
        Inches(4), Inches(2.5)
    )
    col1_box.fill.solid()
    col1_box.fill.fore_color.rgb = LIGHT_GRAY
    col1_box.line.color.rgb = CL_GREEN
    col1_box.line.width = Pt(2)
    
    col1_title = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(3.6), Inches(0.4))
    tf = col1_title.text_frame
    tf.text = "Operational Efficiency"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = CL_DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    y = 2.1
    add_icon_bullet(slide, 1, y, "⚡", "80% faster processing", 1)
    y += 0.5
    add_icon_bullet(slide, 1, y, "💰", "60% cost reduction", 1)
    y += 0.5
    add_icon_bullet(slide, 1, y, "🤖", "70% auto-triage", 1)
    
    # Right column - Customer
    col2_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.3),
        Inches(4), Inches(2.5)
    )
    col2_box.fill.solid()
    col2_box.fill.fore_color.rgb = LIGHT_GRAY
    col2_box.line.color.rgb = CL_GREEN
    col2_box.line.width = Pt(2)
    
    col2_title = slide.shapes.add_textbox(Inches(5.4), Inches(1.5), Inches(3.6), Inches(0.4))
    tf = col2_title.text_frame
    tf.text = "Customer Experience"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = CL_DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    y = 2.1
    add_icon_bullet(slide, 5.4, y, "🕐", "24/7 submission", 1)
    y += 0.5
    add_icon_bullet(slide, 5.4, y, "📱", "Real-time updates", 1)
    y += 0.5
    add_icon_bullet(slide, 5.4, y, "📲", "Mobile-responsive", 1)
    
    # Bottom highlight
    highlight = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(4.2),
        Inches(8.4), Inches(2.3)
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = CL_GREEN
    highlight.line.fill.background()
    
    roi_text = slide.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(7.6), Inches(1.5))
    tf = roi_text.text_frame
    tf.text = "ROI: $1.84M Annual Savings\n58% Cost Reduction"
    for para in tf.paragraphs:
        para.font.size = Pt(32)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    
    # ===== SLIDE 10: Deployment Strategy =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Deployment Strategy")
    add_footer(slide, 10)
    
    # Timeline
    phases = [
        ("Phase 1", "MVP (Current)", ["Core claim submission", "Local deployment", "Docker Compose"]),
        ("Phase 2", "Q3 2026", ["AWS integration", "Kubernetes/EKS", "Production ready"]),
        ("Phase 3", "Q4 2026", ["AI/ML fraud detection", "Mobile apps", "Analytics dashboard"])
    ]
    
    y_start = 1.5
    for i, (phase, timeline, items) in enumerate(phases):
        y = y_start + i * 1.7
        
        # Phase box
        phase_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y),
            Inches(8.4), Inches(1.4)
        )
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = LIGHT_GRAY if i > 0 else CL_GREEN
        phase_box.line.color.rgb = CL_GREEN
        phase_box.line.width = Pt(2)
        
        # Phase title
        title_color = WHITE if i == 0 else CL_DARK_GREEN
        phase_title = slide.shapes.add_textbox(Inches(1), Inches(y + 0.1), Inches(2), Inches(0.3))
        tf = phase_title.text_frame
        tf.text = f"{phase}: {timeline}"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = title_color
        
        # Items
        items_text = " • ".join(items)
        items_box = slide.shapes.add_textbox(Inches(1), Inches(y + 0.5), Inches(7.6), Inches(0.7))
        tf = items_box.text_frame
        tf.text = items_text
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE if i == 0 else DARK_GRAY
    
    # ===== SLIDE 11: Cost Analysis =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Cost Analysis")
    add_footer(slide, 11)
    
    # Current vs New comparison
    # Current System
    current_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.5),
        Inches(4), Inches(4)
    )
    current_box.fill.solid()
    current_box.fill.fore_color.rgb = RGBColor(255, 235, 235)
    current_box.line.color.rgb = RGBColor(220, 53, 69)
    current_box.line.width = Pt(3)
    
    current_title = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(3.6), Inches(0.4))
    tf = current_title.text_frame
    tf.text = "Current System"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 53, 69)
    p.alignment = PP_ALIGN.CENTER
    
    y = 2.3
    items = [
        ("Manual processing", "$2.4M"),
        ("Infrastructure", "$800K"),
        ("Total Annual", "$3.2M")
    ]
    for label, value in items:
        label_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(2), Inches(0.3))
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        
        value_box = slide.shapes.add_textbox(Inches(3.2), Inches(y), Inches(1.2), Inches(0.3))
        tf = value_box.text_frame
        tf.text = value
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(220, 53, 69)
        p.alignment = PP_ALIGN.RIGHT
        y += 0.5
    
    # New System
    new_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.5),
        Inches(4), Inches(4)
    )
    new_box.fill.solid()
    new_box.fill.fore_color.rgb = RGBColor(230, 247, 237)
    new_box.line.color.rgb = CL_GREEN
    new_box.line.width = Pt(3)
    
    new_title = slide.shapes.add_textbox(Inches(5.4), Inches(1.7), Inches(3.6), Inches(0.4))
    tf = new_title.text_frame
    tf.text = "New System"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CL_DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    y = 2.3
    items = [
        ("Automated processing", "$960K"),
        ("Cloud infrastructure", "$400K"),
        ("Total Annual", "$1.36M")
    ]
    for label, value in items:
        label_box = slide.shapes.add_textbox(Inches(5.6), Inches(y), Inches(2), Inches(0.3))
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        
        value_box = slide.shapes.add_textbox(Inches(7.6), Inches(y), Inches(1.2), Inches(0.3))
        tf = value_box.text_frame
        tf.text = value
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = CL_DARK_GREEN
        p.alignment = PP_ALIGN.RIGHT
        y += 0.5
    
    # Savings highlight
    savings_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(5.8),
        Inches(6), Inches(0.8)
    )
    savings_box.fill.solid()
    savings_box.fill.fore_color.rgb = CL_GREEN
    savings_box.line.fill.background()
    
    savings_text = slide.shapes.add_textbox(Inches(2.2), Inches(5.95), Inches(5.6), Inches(0.5))
    tf = savings_text.text_frame
    tf.text = "Annual Savings: $1.84M (58% Reduction)"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # ===== SLIDE 12: Risk Mitigation =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Risk Mitigation")
    add_footer(slide, 12)
    
    risks = [
        ("Technical", [
            "Microservices complexity → Comprehensive monitoring",
            "Data migration → Phased rollout with parallel systems"
        ]),
        ("Operational", [
            "Staff training → 4-week training program",
            "Change management → Dedicated support team"
        ]),
        ("Security", [
            "Data breaches → Multi-layer security, encryption",
            "Compliance → Built-in audit trails, regular audits"
        ])
    ]
    
    y = 1.3
    for category, items in risks:
        # Category header
        cat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y),
            Inches(8.4), Inches(0.5)
        )
        cat_box.fill.solid()
        cat_box.fill.fore_color.rgb = CL_GREEN
        cat_box.line.fill.background()
        
        cat_text = slide.shapes.add_textbox(Inches(1), Inches(y + 0.05), Inches(8), Inches(0.4))
        tf = cat_text.text_frame
        tf.text = f"{category} Risks"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        y += 0.6
        
        # Items
        for item in items:
            item_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(7.6), Inches(0.35))
            tf = item_box.text_frame
            tf.text = f"• {item}"
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            y += 0.45
        
        y += 0.2
    
    # ===== SLIDE 13: Next Steps =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_canada_life_header(slide, "Next Steps")
    add_footer(slide, 13)
    
    # Action items with timeline
    actions = [
        ("Week 1-2", "Immediate Actions", [
            "Executive approval and budget allocation",
            "Form project team and assign roles"
        ]),
        ("Month 1-3", "Short Term", [
            "AWS account setup and infrastructure",
            "Integration with existing systems",
            "User acceptance testing"
        ]),
        ("Month 4-6", "Medium Term", [
            "Pilot launch with select policies",
            "Full production rollout"
        ])
    ]
    
    y = 1.3
    for timeline, phase, items in actions:
        # Timeline badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y),
            Inches(1.5), Inches(0.4)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = CL_GREEN
        badge.line.fill.background()
        
        badge_text = slide.shapes.add_textbox(Inches(0.9), Inches(y + 0.05), Inches(1.3), Inches(0.3))
        tf = badge_text.text_frame
        tf.text = timeline
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Phase title
        phase_text = slide.shapes.add_textbox(Inches(2.5), Inches(y), Inches(3), Inches(0.4))
        tf = phase_text.text_frame
        tf.text = phase
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = CL_DARK_GREEN
        
        y += 0.5
        
        # Items
        for item in items:
            add_icon_bullet(slide, 1, y, "✓", item, 1)
            y += 0.4
        
        y += 0.3
    
    # Call to action
    cta_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(5.8),
        Inches(7), Inches(0.8)
    )
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = CL_GREEN
    cta_box.line.fill.background()
    
    cta_text = slide.shapes.add_textbox(Inches(1.7), Inches(5.95), Inches(6.6), Inches(0.5))
    tf = cta_text.text_frame
    tf.text = "Ready to Transform Claims Processing"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Save
    prs.save('presentation/Canada_Life_Executive_Presentation.pptx')
    print("✅ Styled PowerPoint presentation created: presentation/Canada_Life_Executive_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()

# Made with Bob
